import pptx
import re
from replace import paragraph_process

file = 'ppt_test - 副本.pptx'
ppt = pptx.Presentation(file)
for slide in ppt.slides:
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            print(type(run))
                        if re.search(r'\{(.+)}', paragraph.text):
                            init_run: list = []
                            runs_to_clear: list = []
                            text: list[str] = []
                            left_count: list[int] = []
                            for run in paragraph.runs:
                                # 因为ppt中的run切分的相当细，必须进行一些合并
                                # 由于可能出现嵌套，用{计数标记当前run是否在{}内，}匹配完成则退出
                                if not left_count or left_count[-1] <= 0:
                                    if '{' in run.text:
                                        init_run.append(run)
                                        text.append(run.text)
                                        left_count.append(run.text.count('{'))
                                        left_count[-1] -= run.text.count('}')
                                else:
                                    text[-1] += run.text
                                    runs_to_clear.append(run)
                                    left_count[-1] += run.text.count('{')
                                    left_count[-1] -= run.text.count('}')
                            for run in runs_to_clear:
                                run.text = ''
                            for idx, run in enumerate(init_run):
                                print(text[idx])
                                run.text = re.sub(r'\{(.+?)}', 'test', text[idx])
                                # run.text = re.sub(r'\{(.+)}', func, run.text)
        # if placeholder.text:
        #   placeholder.text = 'test'
ppt.save('test_temp.pptx')
