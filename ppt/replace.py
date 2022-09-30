# 替换ppt中的文本
import functools
import os
import re
from datetime import datetime
from typing import Callable, Optional, TypeVar, Tuple, Type
import downloader
import CVSE_Data
from CVSE_Data import rank_trans, _input
import RankData
import VariablesClass
import pptx

from RankData import Slide, Run, Value, Class, ValidMember, Member, TableType, RootType

#os.chdir('../')
alias_dict: dict[str, str] = {}
with open('数据别名.txt', 'r', encoding='utf-8') as f:
    for line in f.readlines():
        line = line.strip()
        if not line:
            continue
        if line[0] == '#':
            continue
        line = line.split('#')[0]
        if line:
            line_list = [i.strip() for i in line.split('=')]
            if len(line_list) != 2:
                print(f'Invalid line: {line_list}')
                continue
            alias_dict[line_list[0]] = line_list[1]

T = TypeVar('T', int, float, str, datetime)


def load_data(index: int, rank: int) -> RankData.RankData:
    default_dir: str = f'../{rank_trans[rank]}_{index}'
    default_file: list[str] = [f'{rank_trans[rank]}_{index}_save_backup.csv', f'{rank_trans[rank]}_{index}.xlsx',
                               f'{index}.xlsx']
    required_keys: list[str] = ['播放增量', '弹幕增量', '评论增量', '收藏增量', '硬币增量', '分享增量', '点赞增量']
    if os.path.exists(default_dir):
        for file in default_file:
            if os.path.exists(os.path.join(default_dir, file)):
                print(f'正在读取{file}')
                video_data: list[CVSE_Data.Data] = CVSE_Data.read(os.path.join(default_dir, file),
                                                                  required_keys=required_keys)
                for i in filter(lambda x: x['收录'] in ['', '1'], video_data):
                    i['收录'] = 1
                return RankData.RankData(video_data, index=index, rank=rank)
    else:
        os.mkdir(default_dir)
    _, end_time = RankData.calculate_time(rank, index)
    print(f'尝试下载第{index}期数据')
    if downloader.download_history_data(end_time, rank, index, default_dir):
        file = f'{index}.xlsx'
        print(f'正在读取{file}')
        video_data = CVSE_Data.read(os.path.join(default_dir, file), required_keys=required_keys)
        for i in filter(lambda x: x['收录'] == '', video_data):
            i['收录'] = 1
        return RankData.RankData(video_data, index=index, rank=rank)
    else:
        print(f'未在{default_dir}找到{"或".join(default_file)}，且下载失败，请手动下载')
        input('按任意键退出')
        raise FileNotFoundError


def get_data(index: int, rank: int) -> Tuple[dict[str, Class], dict[str, TableType]]:
    # 读取近十期数据,dict的key为期数，value为Rank_data
    data_dict: dict[int, RankData.RankData] = {}
    last_dur = min(10, index)  # 读取这么多期过往数据
    for i in range(last_dur):
        data_dict[index - i] = load_data(index - i, rank)
    history_index: int = RankData.calculate_history(rank, index)
    history_data: Optional[RankData.RankData] = load_data(history_index, rank) if history_index else None
    if any(i is None for i in data_dict.values()):
        print()
    Root: dict[str, Class] = {
        'pres': data_dict[index],
        'prev': None,
        'history': history_data,
        'pres_delta': RankData.RankDataDelta(pres_data=data_dict[index],
                                             prev_data=data_dict[index - 1]) if index > 1 else None,
        'pres_rate': RankData.RankDataRate(pres_data=data_dict[index],
                                           prev_data=data_dict[index - 1]) if index > 1 else None,

    }
    Table_value: dict[str, TableType] = {
        'top1view_recent_10_view': [data_dict[index - i]['top1_view']['播放增量'] for i in range(last_dur)],
        'top1view_recent_10_title': [data_dict[index - i]['top1_view']['标题'] for i in range(last_dur)],
        'total_view_recent_10': [data_dict[index - i]['view'] for i in range(last_dur)],
        'startpt_recent_10': [data_dict[index - i]['startpt'] for i in range(last_dur)],
        'side_startpt_recent10': [data_dict[index - i]['side_startpt'] for i in range(last_dur)],
    }
    if last_dur != 1:
        Root['prev'] = data_dict[index - 1]
    return Root, Table_value


def data_lookup(root: RootType, *args: str) -> ValidMember:
    if root is None:
        raise ValueError('root is None')
    if len(args) == 1:
        # print(args)
        if args[0] in root:
            if root[args[0]] is None:
                raise ValueError(f'{args[0]} is None')
            return root[args[0]]
        else:
            print(f'Invalid key: {args[0]}')
            # input()
            raise ValueError
    else:
        if args[0] in root:
            return data_lookup(root[args[0]], *args[1:])
        else:
            print(f'Invalid key: {args[0]}')
            # input()
            raise ValueError


def type_split(name: str) -> tuple[str, Optional[Type[Value]]]:
    # 由输入的字符串（类型名），返回解析到的类型名和类型类
    if any(i in name for i in ['int', 'float']):
        for valuetype in ['int', 'float']:
            if name == valuetype:
                return valuetype, VariablesClass.IntValue if valuetype == 'int' else VariablesClass.FloatValue
            else:
                result = re.search(r'color' + valuetype, name)
                if result:
                    color_index: int = 0
                    if result.group(0) != name:
                        result = re.search(r'color' + valuetype + r'_\$(\d+)', name)
                        if result and result.group(0) == name:
                            color_index = int(result.group(1))
                        else:
                            raise ValueError(f'Invalid type: {name}')
                    return valuetype, VariablesClass.get_ColorValue_class(color_index, valuetype)
    if name == 'str':
        return name, str
    if name == 'time':
        return name, VariablesClass.Time
    if name == 'pic':
        return name, VariablesClass.PicturePlacer
    if name in ['Table']:
        return name, None   # WIP
    raise ValueError(f'Invalid type: {name}')


def format_func(arg: str, root: RootType) \
        -> Tuple[Value, Optional[Callable[[RankData.Shapes], None]]]:
    # 格式化成功返回相应对象，否则返回原字符串,第二个参数是对ppt的操作，带一个参数slide类
    if arg in alias_dict:
        name = alias_dict[arg.strip()]
    else:
        name = arg.strip()
    args_list: list[str] = re.split(r'\s(?![^(]*\))', name) # 匹配不在括号内的空格,感谢
    if len(args_list) == 1:
        return arg, None
    try:
        value_type, value_class = type_split(args_list[0])
    except ValueError as e:
        print(f'{e} in {name}')
        return '{' + arg + '}', None
    if value_class is not None:
        try:
            if args_list[1].endswith('"') and args_list[1].startswith('"'):
                value: Value = value_class(args_list[1], *args_list[2:])
            else:
                value = value_class(data_lookup(root, *(args_list[1].split(':'))), *args_list[2:])
            operation: Optional[Callable[[Slide], None]] \
                = value.get_callable() if isinstance(value, VariablesClass.BaseValue) else None
            return value, operation
        except ValueError as e:
            print(e)
            print(f'Invalid data or args: {args_list[1:]}')
            return '{' + arg + '}', None
    else:
        return '{' + arg + '}', None


def action(match: str, root: RootType, slide: Slide, run: Run) -> str:
    # 格式化字符串并对幻灯片执行操作
    arg: str = match
    result, operation = format_func(arg, root)
    if operation is not None:
        operation(slide.shapes)
    if isinstance(result, VariablesClass.BaseValue):
        if result.with_color():
            run.font.color.rgb = pptx.dml.color.RGBColor.from_string(result.get_color())
    try:
        text = str(result)
    except ValueError as e:
        print(f'{arg}格式化失败')
        print(e)
        text = '{' + arg + '}'
    text = re.sub(r'\\s', ' ', text)  # 替换空格
    return text


def replace_text(ppt, func: Callable[[str, RootType, Slide, Run], str], root: RootType):
    # 替换幻灯片中的文本
    for slide in ppt.slides:
        # print(type(slide))
        func = functools.partial(func, root=root, slide=slide)  # 将当前slide传入func，之后可以在func内部对slide进行操作
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # print(paragraph.text)
                    paragraph_process(func, paragraph)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph_process(func, paragraph)


def paragraph_process(func: Callable[[str, Run], str], paragraph):
    # 替换段落中的文本
    if re.search(r'\{(.+)}', paragraph.text):
        print(paragraph.text)
        init_run: list = []  # 每段开头的run，这一段的格式全部按照开头run的格式
        runs_to_clear: list = []  # 需要清楚文本的后继run，这些文本暂存在text中，之后复制到开头的run里
        text: list[str] = []
        left_count: list[int] = []  # 统计{的数量，标记当前是否在{}内
        for run in paragraph.runs:
            if not left_count or left_count[-1] <= 0:
                if '{' in run.text:
                    init_run.append(run)
                    text.append(run.text)
                    left_count.append(run.text.count('{'))
                    left_count[-1] -= run.text.count('}')
            else:
                text[-1] += run.text
                runs_to_clear.append(run)
                left_count[-1] += run.text.count('{')
                left_count[-1] -= run.text.count('}')
        for run in runs_to_clear:
            run.text = ''
        for idx, run in enumerate(init_run):
            # (text[idx])
            run.text = text[idx]
            func = functools.partial(func, run=run)
            _left_count = run.text.count('{')
            _right_count = run.text.count('}')
            if _left_count != _right_count:
                print(f'Invalid format: {run.text}')
                continue
            stack: list[int] = []  # 栈，用于匹配大括号
            ori_text: str = run.text
            i: int = 0
            while i < len(run.text):
                if run.text[i] == '{':
                    stack.append(i)
                if run.text[i] == '}':
                    if not stack:
                        print(f'Invalid format: {run.text}')
                        break
                    left = stack.pop()
                    left_str, replaced_str, right_str = run.text[:left], run.text[left + 1:i], run.text[i + 1:]
                    # 注意这里的replaced_str是不包含大括号的,i+1越界时会返回空字符串不需要特别判断
                    print(left_str, replaced_str, right_str, sep='|')
                    # 注意left和i的位置上都是花括号，可以直接舍弃
                    new_str = func(replaced_str)
                    if new_str == '{' + replaced_str + '}':
                        # 这里是因为如果解析失败，会返回原字符串两边加上花括号。只返回原字符串的话，并不是匹配失败
                        print(f'Invalid format: {run.text}')
                        break
                    run.text = left_str + new_str + right_str
                    i = left + len(new_str) - 1
                i += 1
            if stack:
                print(f'Invalid format: {ori_text}')
                run.text = ori_text
            # run.text, flag = re.subn(r'(?!\\)\{(.+?)(?!\\)}', func, text[idx], flag)
            # run.text = re.sub(r'\{(.+)}', func, run.text)


def init() -> Tuple[int, int, dict[str, Class], dict[str, TableType]]:
    _rank = input("请输入国产榜/SV/utau，1=国产榜，2=SV刊，3=UTAU刊\n")
    while _rank not in ['1', '2']:
        print('输入错误')
        _rank = input("请输入国产榜/SV/utau，1=国产榜，2=SV刊，3=UTAU刊\n")
    rank = int(_rank) - 1
    index = _input("请输入待处理排行榜期数，如 133 50\n", lambda x: x.isdigit())
    while not index.isdigit():
        print('输入的不是数字')
        index = input("请输入待处理排行榜期数，如 133 50\n")
    index = int(index)
    root, table_root = get_data(index, rank)
    return index, rank, root, table_root


# print(format_func('int 123.456 ,').with_color())
# ppt.save('test.pptx')
_index, _rank, _root, _table_root = init()
file: str = 'ppt_demo.pptx'
ppt = pptx.Presentation(file)
replace_text(ppt, action, _root)
ppt.save('test.pptx')
