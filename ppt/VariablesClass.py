import os
from abc import ABCMeta
from datetime import datetime
from functools import lru_cache, partial
from io import BytesIO
from typing import Optional, TypeVar, Callable, Literal
import pptx
from PIL import Image
import re

_init_flag = False


@lru_cache
def load_color_map() -> dict[int, tuple[str, str]]:
    color_map: dict[int, tuple[str, str]] = {}
    print(1)
    with open('color.txt', 'r', encoding='utf-8') as f:
        for line in f.readlines():
            _line = line.strip()
            if _line[0] == '#':
                continue
            valid_text = _line.split('#')[0]
            if valid_text:
                __line = [i.strip() for i in valid_text.split('=')]
                temp_list = [i.strip() for i in __line[1].split(',')]
                if len(__line) != 2 or len(temp_list) != 2:
                    print(f'Invalid line: {__line}')
                    continue
                color_map[int(__line[0])] = temp_list[0], temp_list[1]
    return color_map


# print(123)

T = TypeVar('T', int, float, datetime)

class BaseValue(metaclass=ABCMeta):
    # 基础值类，所有值类的基类
    def __init__(self, *args, **kwargs):
        self.value: Optional[T] = None
        self.format_list: Optional[list[str]] = None

    def with_color(self) -> bool:
        return False

    def is_positive(self) -> bool:
        if self.value is None:
            raise NotImplementedError
        return self.value >= 0

    def get_color(self) -> str:
        print(f'Invalid call: {self.__class__.__name__}.get_color()')
        return ''

    def __str__(self) -> str:
        if not self.format_list:
            return f'{self.value}'
        else:
            if isinstance(self.value, int):
                mark_order: list[str] = ['fill', 'align', r'\+', 'width', r'\,']
            elif isinstance(self.value, float):
                mark_order = ['fill', 'align', r'\+', r'\,', r'\.\d', 'width', r'\.\d%']
            else:
                # 这里只实现了int和float的格式化，其他类型的格式化需要自己实现
                raise NotImplementedError
            valid_list: list[str] = []
            format_str: str = '{:'
            align_flag: bool = False
            align_arg: str = ''
            for i in self.format_list:
                result = re.match(r'(\S)([<>])(\d+)', i)
                # 由于顺序问题，这里必须先对对齐方式进行处理，将其拆散，才能符合python的格式化语法
                if result:
                    if align_flag or result.group(0) != i:
                        # 防止输入重复格式
                        raise ValueError('Duplicate align format')
                    else:
                        align_flag = True
                        align_arg = i
                        self.format_list.append(result.group(1))
                        self.format_list.append(result.group(2))
                        self.format_list.append(result.group(3))
                        mark_order[mark_order.index('fill')] = result.group(1)
                        mark_order[mark_order.index('align')] = result.group(2)
                        mark_order[mark_order.index('width')] = result.group(3)
            if align_flag:
                self.format_list.remove(align_arg)
            for i in mark_order:
                for j in self.format_list:
                    result = re.search(i, j)
                    if result and result.group(0) == j:
                        format_str += j
                        valid_list.append(j)
            if set(valid_list) != set(self.format_list):
                raise ValueError(
                    f'Invalid format: {self.format_list} for Value: {self.value} of type {type(self.value)}')
            if isinstance(self.value, int):
                format_str += 'd}'
            elif isinstance(self.value, float):
                if any(re.search(r'\.\d%', i) for i in self.format_list):
                    format_str += '}'
                else:
                    format_str += 'f}'
            return format_str.format(self.value)

    def get_callable(self) -> Optional[Callable]:
        return None


class BaseColor(BaseValue, metaclass=ABCMeta):
    # 有颜色的值
    def __init__(self, *args, **kwargs):
        super(BaseColor, self).__init__(*args, **kwargs)
        self.color_set: Optional[tuple] = None

    def get_color(self):
        return self.color_set[0] if self.is_positive() else self.color_set[1]

    def with_color(self):
        return True


class IntValue(BaseValue):
    def __init__(self, value, *args: str, **kwargs):
        super().__init__()
        self.value: int = int(value)
        self.format_list: list[str] = list(args)


class FloatValue(BaseValue):
    def __init__(self, value, *args: str, **kwargs):
        super().__init__()
        self.value: float = float(value)
        self.format_list: list[str] = list(args)


@lru_cache  # 否则每次都会重新创建一个新的类
def get_ColorValue_class(index: int, value_type: str) -> type[BaseColor]:
    color_set = load_color_map()[index]
    if value_type == 'int':
        class ColorValueInt(IntValue, BaseColor):
            def __init__(self, value, *args: str, **kwargs):
                super().__init__(value, *args, **kwargs)
                self.color_set: tuple[str, str] = color_set

            def with_color(self):
                return True

        return ColorValueInt
    elif value_type == 'float':
        class ColorValueFloat(FloatValue, BaseColor):
            def __init__(self, value, *args: str, **kwargs) -> None:
                super().__init__(value, *args, **kwargs)
                self.color_set: tuple[str, str] = color_set

            def with_color(self) -> Literal[True]:
                return True

        return ColorValueFloat
    else:
        raise ValueError(f'Invalid value type: {value_type}')
    # print(ColorValue.__mro__)


class Time(BaseValue):
    def __init__(self, value: datetime, time_format: str, *args, **kwargs):
        super().__init__()
        self.value: datetime = value
        try:
            self.time_format: str = self.value.strftime(time_format)
        except ValueError:
            raise ValueError(f'Invalid time format: {time_format}')

    def __str__(self) -> str:
        return self.time_format

    def is_positive(self) -> bool:
        print('时间量没有正负之分')
        return True


class PicturePlacer(BaseValue):
    def __init__(self, pic: Image.Image | str, *args: list[str], **kwargs):
        super().__init__()
        self.x: float
        self.y: float
        self.width: float
        self.height: float
        self.pic_file: str | BytesIO
        if len(args) != 2:
            raise ValueError('Invalid picture args')
        try:
            self.x, self.y = eval(args[0])
            self.width, self.height = eval(args[1])
        except Exception as e:
            raise ValueError(f'Invalid picture position or size: {e}')
        if isinstance(pic, Image.Image):
            bytes_io = BytesIO()
            pic.save(bytes_io, format='JPEG')
            self.pic_file = bytes_io
        else:
            self.pic_file = pic

    def __str__(self):
        return ''

    def get_callable(self) -> Optional[Callable[[pptx.shapes.shapetree.SlideShapes], None]]:
        x = pptx.util.Cm(self.x)
        y = pptx.util.Cm(self.y)
        width = pptx.util.Cm(self.width)
        height = pptx.util.Cm(self.height)
        return partial(pptx.shapes.shapetree.SlideShapes.add_picture, image_file=self.pic_file, left=x, top=y, width=width, height=height)




# func = get_ColorValue_class(1, 'float')
# print(func(1232.66575, '+', ',', '.5%'))
# print(ColorValue(123234, ('red', 'green'), '+', ','))
